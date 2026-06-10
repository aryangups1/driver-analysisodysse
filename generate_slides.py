"""
generate_slides.py
Run from project root:  python generate_slides.py
Produces:  odysse_management_deck.pptx
"""

import sys, os, io, warnings
warnings.filterwarnings("ignore")
sys.path.insert(0, os.path.dirname(__file__))

import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
import plotly.io as pio

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import db
from zones import parse_dms, assign_zone
from config import TOP_DRIVER_IDS, BAD_DRIVER_IDS, DRIVER_NAMES

# ── Constants ────────────────────────────────────────────────────────────────
_WEST_LON  = -0.12
W, H       = Inches(13.33), Inches(7.5)   # 16:9

# Colours
C_BG       = RGBColor(0x0f, 0x17, 0x2a)   # slide background
C_CARD     = RGBColor(0x1e, 0x29, 0x3b)   # card bg
C_WHITE    = RGBColor(0xff, 0xff, 0xff)
C_GREEN    = RGBColor(0x22, 0xc5, 0x5e)
C_RED      = RGBColor(0xef, 0x44, 0x44)
C_AMBER    = RGBColor(0xf5, 0x9e, 0x0b)
C_BLUE     = RGBColor(0x3b, 0x82, 0xf6)
C_MUTED    = RGBColor(0x94, 0xa3, 0xb8)
C_ACCENT   = RGBColor(0x38, 0xbd, 0xe8)

ALL_IDS    = TOP_DRIVER_IDS + BAD_DRIVER_IDS

# ── Plotly theme ─────────────────────────────────────────────────────────────
_PLOTLY_BASE = dict(
    paper_bgcolor="#0f172a",
    plot_bgcolor="#1e293b",
    font=dict(color="#e2e8f0", size=13),
    margin=dict(l=40, r=20, t=50, b=40),
)

def _fig_to_img(fig, w=800, h=440):
    fig.update_layout(**_PLOTLY_BASE)
    return io.BytesIO(pio.to_image(fig, format="png", width=w, height=h, scale=3))

# ── PPTX helpers ─────────────────────────────────────────────────────────────
def _new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def _blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    slide  = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return slide

def _add_text(slide, text, left, top, width, height, size=14,
              bold=False, color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def _add_image(slide, img_bytes, left, top, width, height):
    img_bytes.seek(0)
    slide.shapes.add_picture(img_bytes, left, top, width, height)

def _add_rect(slide, left, top, width, height, fill=C_CARD, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def _slide_header(slide, title, subtitle=None):
    _add_text(slide, title, Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.7),
              size=28, bold=True, color=C_WHITE)
    if subtitle:
        _add_text(slide, subtitle, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4),
                  size=13, color=C_MUTED, italic=True)

def _metric_card(slide, label, value, left, top, w=Inches(2.4), h=Inches(1.1),
                 val_color=C_WHITE, delta=None, delta_color=C_GREEN):
    _add_rect(slide, left, top, w, h, fill=C_CARD)
    _add_text(slide, label, left + Inches(0.12), top + Inches(0.08),
              w - Inches(0.2), Inches(0.35), size=10, color=C_MUTED)
    _add_text(slide, value, left + Inches(0.12), top + Inches(0.38),
              w - Inches(0.2), Inches(0.5), size=22, bold=True, color=val_color)
    if delta:
        _add_text(slide, delta, left + Inches(0.12), top + Inches(0.78),
                  w - Inches(0.2), Inches(0.25), size=10, color=delta_color)

def _callout(slide, text, left, top, width, height, border_color=C_AMBER):
    _add_rect(slide, left, top, width, height, fill=RGBColor(0x1e,0x1e,0x2e), line=border_color)
    _add_text(slide, text, left + Inches(0.15), top + Inches(0.12),
              width - Inches(0.3), height - Inches(0.2),
              size=12, color=C_WHITE)

# ── Data loading ─────────────────────────────────────────────────────────────
print("Loading data from DB…")

perf_top  = db.load_comparison_performance(TOP_DRIVER_IDS)
perf_bad  = db.load_comparison_performance(BAD_DRIVER_IDS)

flow_top  = db.load_comparison_flow(TOP_DRIVER_IDS)
flow_bad  = db.load_comparison_flow(BAD_DRIVER_IDS)

trips_top = db.load_comparison_trips(TOP_DRIVER_IDS)
trips_bad = db.load_comparison_trips(BAD_DRIVER_IDS)

coords    = db.load_all_driver_coords(sample_per_driver=80, days_back=60)

gaps_top     = db.load_gap_accepted(TOP_DRIVER_IDS, days_back=60)
gaps_bad     = db.load_gap_accepted(BAD_DRIVER_IDS, days_back=60)
flow_yousuf  = db.load_comparison_flow([219])   # outlier spotlight — not in TOP_DRIVER_IDS

print("Data loaded. Enriching zones…")

# ── Zone enrichment helper ────────────────────────────────────────────────────
def _enrich(df, lat_col="pickup_lat_long", lon_col=None):
    coords_parsed = df[lat_col].apply(lambda x: parse_dms(str(x or "")))
    df = df.copy()
    df["plat"] = [c[0] for c in coords_parsed]
    df["plon"] = [c[1] for c in coords_parsed]
    df = df.dropna(subset=["plat","plon"])
    df["pickup_zone"] = [assign_zone(a, b) for a, b in zip(df["plat"], df["plon"])]
    df["is_west"]     = df["plon"] < _WEST_LON
    if lon_col and lon_col in df.columns:
        dcoords = df[lon_col].apply(lambda x: parse_dms(str(x or "")))
        df["dlat"] = [c[0] for c in dcoords]
        df["dlon"] = [c[1] for c in dcoords]
        df = df.dropna(subset=["dlat","dlon"])
        df["dropoff_zone"] = [assign_zone(a, b) for a, b in zip(df["dlat"], df["dlon"])]
    return df

acc_top = _enrich(
    flow_top[flow_top["status"].isin(["completed","Finished"])].copy(),
    lon_col="dropoff_latlong"
)
acc_bad = _enrich(
    flow_bad[flow_bad["status"].isin(["completed","Finished"])].copy(),
    lon_col="dropoff_latlong"
)
acc_top["fare"] = pd.to_numeric(acc_top["trip_price_in_pound"], errors="coerce")
acc_bad["fare"] = pd.to_numeric(acc_bad["trip_price_in_pound"], errors="coerce")

acc_yousuf = _enrich(
    flow_yousuf[flow_yousuf["status"].isin(["completed","Finished"])].copy(),
    lon_col="dropoff_latlong"
)
acc_yousuf["fare"] = pd.to_numeric(acc_yousuf["trip_price_in_pound"], errors="coerce")

# Inner / outer sub-region classification (applied to all accepted trips)
def _subregion(plon):
    if plon < -0.25:   return "Outer West"
    if plon < -0.12:   return "Inner West"
    if plon <  0.0:    return "Inner East"
    return "Outer East"

acc_top["subregion"] = acc_top["plon"].apply(_subregion)
acc_bad["subregion"] = acc_bad["plon"].apply(_subregion)

# Fleet map coords
cat_map = {}  # driver_id → category
try:
    cat_df = pd.read_csv(os.path.join(os.path.dirname(__file__), "driver_categories.csv"))
    cat_map = dict(zip(cat_df["dim_driver_id"].astype(int), cat_df["category"].astype(str)))
except Exception:
    pass

coords_df = pd.DataFrame(coords, columns=["dim_driver_id","driver_full_name","pickup_lat_long"])
coords_df = coords_df.dropna(subset=["dim_driver_id"])
coords_df["dim_driver_id"] = coords_df["dim_driver_id"].astype(float).astype(int)
coords_df["category"] = coords_df["dim_driver_id"].map(cat_map).fillna("Unknown")
_pu = coords_df["pickup_lat_long"].apply(lambda x: parse_dms(str(x or "")))
coords_df["plat"] = [c[0] for c in _pu]
coords_df["plon"] = [c[1] for c in _pu]
coords_df = coords_df.dropna(subset=["plat","plon"])
coords_df = coords_df[(coords_df["plat"].between(51.2,51.8)) & (coords_df["plon"].between(-0.6,0.3))]

# Gap computation helper
def _compute_gaps(trips_df):
    df = trips_df.copy()
    df["pickup_dt"]  = pd.to_datetime(df["pickedup_trip_datetime"], errors="coerce")
    df["dropoff_dt"] = pd.to_datetime(df["dropoff_trip_datetime"],  errors="coerce")
    df = df.dropna(subset=["pickup_dt","dropoff_dt"]).sort_values(["dim_driver_id","pickup_dt"])
    df["next_pickup"] = df.groupby("dim_driver_id")["pickup_dt"].shift(-1)
    df["gap_min"] = ((df["next_pickup"] - df["dropoff_dt"]).dt.total_seconds() / 60).clip(0, 180)
    return df.dropna(subset=["gap_min"])

gaps_top_df = _compute_gaps(pd.DataFrame(gaps_top,
    columns=["dim_driver_id","driver_full_name","pickup_lat_long","dropoff_latlong",
             "trip_price_in_pound","distance_in_miles","pob_duration_in_min",
             "pickup_duration_in_min","trips_hr","source",
             "pickedup_trip_datetime","dropoff_trip_datetime"]))
gaps_bad_df = _compute_gaps(pd.DataFrame(gaps_bad,
    columns=["dim_driver_id","driver_full_name","pickup_lat_long","dropoff_latlong",
             "trip_price_in_pound","distance_in_miles","pob_duration_in_min",
             "pickup_duration_in_min","trips_hr","source",
             "pickedup_trip_datetime","dropoff_trip_datetime"]))

print("Done. Building slides…")

# ── Aggregate metrics ─────────────────────────────────────────────────────────
perf_top_df = pd.DataFrame(perf_top,
    columns=["dim_driver_id","driver_name","total_rides","online_hrs",
             "total_revenue","rph","utilisation","acceptance","avg_fare"])
perf_bad_df = pd.DataFrame(perf_bad,
    columns=["dim_driver_id","driver_name","total_rides","online_hrs",
             "total_revenue","rph","utilisation","acceptance","avg_fare"])

def _avg(df, col): return round(pd.to_numeric(df[col], errors="coerce").mean(), 2)

top_rph    = _avg(perf_top_df, "rph")
bad_rph    = _avg(perf_bad_df, "rph")
top_fare   = _avg(perf_top_df, "avg_fare")
bad_fare   = _avg(perf_bad_df, "avg_fare")
top_acc    = _avg(perf_top_df, "acceptance")
bad_acc    = _avg(perf_bad_df, "acceptance")

west_top   = round(acc_top["is_west"].mean() * 100, 1)
west_bad   = round(acc_bad["is_west"].mean() * 100, 1)
wfare_top  = round(acc_top[acc_top["is_west"]]["fare"].mean(), 2)
efare_top  = round(acc_top[~acc_top["is_west"]]["fare"].mean(), 2)

z3_top     = round((acc_top["pickup_zone"] == 3).mean() * 100, 1)
z3_bad     = round((acc_bad["pickup_zone"] == 3).mean() * 100, 1)
z1_top     = round((acc_top["pickup_zone"] == 1).mean() * 100, 1)
z1_bad     = round((acc_bad["pickup_zone"] == 1).mean() * 100, 1)

chain_top  = round(acc_top[acc_top["pickup_zone"]==3]["dropoff_zone"].eq(3).mean() * 100, 1) \
             if "dropoff_zone" in acc_top.columns and not acc_top.empty else 0
chain_bad  = round(acc_bad[acc_bad["pickup_zone"]==3]["dropoff_zone"].eq(3).mean() * 100, 1) \
             if "dropoff_zone" in acc_bad.columns and not acc_bad.empty else 0

gap_top_avg = round(gaps_top_df["gap_min"].mean(), 1)
gap_bad_avg = round(gaps_bad_df["gap_min"].mean(), 1)

# ── BUILD PRESENTATION ────────────────────────────────────────────────────────
prs = _new_prs()

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_add_text(s, "ODYSSE FLEET", Inches(1), Inches(1.8), Inches(11), Inches(1.2),
          size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
_add_text(s, "PERFORMANCE ANALYSIS", Inches(1), Inches(2.9), Inches(11), Inches(1),
          size=44, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
_add_text(s, "Driver Behaviour · Zone Strategy · Gap Optimisation · Acceptance Rate",
          Inches(1), Inches(4.0), Inches(11), Inches(0.5),
          size=15, color=C_MUTED, align=PP_ALIGN.CENTER)
_add_text(s, "Internal — Management Review   ·   June 2026",
          Inches(1), Inches(6.6), Inches(11), Inches(0.4),
          size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2: FLEET MAP
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Fleet Positioning — Where Are We Operating?",
              "Pickup locations for all active drivers (last 60 days). Colour = performance category.")

cat_order  = ["A","B1","B2","C1","C2","D","Unknown"]
cat_colors = {"A":"#22c55e","B1":"#84cc16","B2":"#a3e635",
              "C1":"#f59e0b","C2":"#fb923c","D":"#ef4444","Unknown":"#64748b"}
fig_map = px.scatter_mapbox(
    coords_df, lat="plat", lon="plon",
    color="category", category_orders={"category": cat_order},
    color_discrete_map=cat_colors,
    hover_data={"driver_full_name": True, "category": True, "plat": False, "plon": False},
    zoom=9.5, center={"lat":51.505,"lon":-0.09},
    mapbox_style="carto-darkmatter",
    height=530, title="",
    opacity=0.8,
)
fig_map.update_traces(marker_size=6)
fig_map.update_layout(paper_bgcolor="#0f172a", margin=dict(l=0,r=0,t=0,b=0),
                      legend=dict(orientation="h", y=0.02, x=0.01,
                                  font=dict(color="#e2e8f0", size=11)))
_add_image(s, _fig_to_img(fig_map, w=1100, h=530), Inches(0.4), Inches(1.15), Inches(8.8), Inches(5.4))

# Legend callout
_add_rect(s, Inches(9.5), Inches(1.2), Inches(3.5), Inches(5.2), fill=C_CARD)
_add_text(s, "CATEGORY KEY", Inches(9.65), Inches(1.3), Inches(3.2), Inches(0.35),
          size=11, bold=True, color=C_MUTED)
cat_desc = [
    ("A",  C_GREEN,              "Elite — top RPH, strategic"),
    ("B1", RGBColor(0x84,0xcc,0x16), "High performers"),
    ("B2", RGBColor(0xa3,0xe6,0x35), "Solid performers"),
    ("C1", C_AMBER,              "Average, improving"),
    ("C2", RGBColor(0xfb,0x92,0x3c), "Below average"),
    ("D",  C_RED,                "Underperforming"),
]
for i, (cat, col, desc) in enumerate(cat_desc):
    y = Inches(1.75) + i * Inches(0.62)
    _add_rect(s, Inches(9.65), y, Inches(0.35), Inches(0.35), fill=col)
    _add_text(s, f"Cat {cat} — {desc}", Inches(10.1), y, Inches(2.8), Inches(0.35),
              size=11, color=C_WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3: THE NUMBER GAP — TOP vs BAD DRIVERS
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "The Performance Gap is Structural, Not Random",
              "Same platform, same city, same hours. These numbers should not be this different.")

_metric_card(s, "Top driver avg RPH",      f"£{top_rph:.2f}", Inches(0.5),  Inches(1.4),
             val_color=C_GREEN)
_metric_card(s, "Bad driver avg RPH",      f"£{bad_rph:.2f}", Inches(3.1),  Inches(1.4),
             val_color=C_RED, delta=f"£{bad_rph - top_rph:.2f} vs top", delta_color=C_RED)
_metric_card(s, "Top avg fare per trip",   f"£{top_fare:.2f}", Inches(5.7),  Inches(1.4),
             val_color=C_GREEN)
_metric_card(s, "Bad avg fare per trip",   f"£{bad_fare:.2f}", Inches(8.3),  Inches(1.4),
             val_color=C_RED, delta=f"£{bad_fare - top_fare:.2f} per trip", delta_color=C_RED)
_metric_card(s, "Top acceptance rate",     f"{top_acc:.0f}%",  Inches(10.9), Inches(1.4),
             val_color=C_GREEN)

# Named driver bar chart: RPH
perf_all  = pd.concat([
    perf_top_df.assign(group="Top drivers"),
    perf_bad_df.assign(group="Comparison drivers"),
], ignore_index=True)
perf_all["driver_name"] = perf_all["dim_driver_id"].astype(int).map(DRIVER_NAMES).fillna(perf_all["driver_name"])
perf_all["rph"]         = pd.to_numeric(perf_all["rph"], errors="coerce")
perf_all["avg_fare"]    = pd.to_numeric(perf_all["avg_fare"], errors="coerce")
perf_all = perf_all.dropna(subset=["rph"]).sort_values("rph", ascending=True)

fig_rph = px.bar(
    perf_all, x="rph", y="driver_name", orientation="h",
    color="group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison drivers": "#ef4444"},
    text="rph",
    labels={"rph": "Revenue per Hour (£)", "driver_name": ""},
    title="RPH by driver — Top vs Comparison",
    height=390,
)
fig_rph.update_traces(texttemplate="£%{text:.2f}", textposition="outside")
fig_rph.add_vline(x=top_rph, line_dash="dash", line_color="#22c55e",
                  annotation_text=f"Top avg £{top_rph:.2f}")
fig_rph.add_vline(x=bad_rph, line_dash="dash", line_color="#ef4444",
                  annotation_text=f"Bad avg £{bad_rph:.2f}", annotation_position="bottom right")
fig_rph.update_layout(showlegend=True, legend=dict(orientation="h",y=1.05),
                      margin=dict(l=200, r=80, t=50, b=40))
_add_image(s, _fig_to_img(fig_rph, w=900, h=430), Inches(0.5), Inches(2.6), Inches(8.1), Inches(4.6))

# Fare comparison bar
perf_fare = perf_all.sort_values("avg_fare", ascending=True)
fig_fare2 = px.bar(
    perf_fare, x="avg_fare", y="driver_name", orientation="h",
    color="group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison drivers": "#ef4444"},
    text="avg_fare",
    labels={"avg_fare": "Avg fare (£)", "driver_name": ""},
    title="Avg fare / trip",
    height=430,
)
fig_fare2.update_traces(texttemplate="£%{text:.2f}", textposition="outside")
fig_fare2.update_layout(showlegend=False, margin=dict(l=200, r=80, t=50, b=40))
_add_image(s, _fig_to_img(fig_fare2, w=520, h=430), Inches(8.8), Inches(2.6), Inches(4.3), Inches(4.6))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4: EAST vs WEST — THE FUNDAMENTAL SPLIT
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "East vs West: The Same City, Two Different Games",
              "West of Charing Cross (−0.12°): premium demand, better fares, cleaner pings. "
              "East: same high-value pings exist — but buried in noise.")

_metric_card(s, "Top drivers — west %",   f"{west_top:.0f}%",  Inches(0.5),  Inches(1.3), val_color=C_GREEN, w=Inches(2.8))
_metric_card(s, "Bad drivers — west %",   f"{west_bad:.0f}%",  Inches(3.6),  Inches(1.3), val_color=C_RED,
             delta=f"{west_bad - west_top:+.0f}% vs top", delta_color=C_RED, w=Inches(2.8))
_metric_card(s, "West avg fare",           f"£{wfare_top:.2f}", Inches(6.7),  Inches(1.3), val_color=C_GREEN, w=Inches(2.8))
_metric_card(s, "East avg fare",           f"£{efare_top:.2f}", Inches(9.8),  Inches(1.3), val_color=C_AMBER,
             delta=f"£{efare_top - wfare_top:.2f} vs west", delta_color=C_RED, w=Inches(2.8))

# Single combined scatter map — top (green) vs bad (red) pickups
_ew_sample = pd.concat([
    acc_top[["plat","plon"]].assign(Group="Top drivers"),
    acc_bad[["plat","plon"]].assign(Group="Comparison"),
], ignore_index=True)
_ew_sample = _ew_sample[_ew_sample["plat"].between(51.3, 51.7) & _ew_sample["plon"].between(-0.55, 0.25)]
# Sample to keep map readable
_ew_top = _ew_sample[_ew_sample["Group"]=="Top drivers"].sample(min(600, len(_ew_sample[_ew_sample["Group"]=="Top drivers"])), random_state=42)
_ew_bad = _ew_sample[_ew_sample["Group"]=="Comparison"].sample(min(600, len(_ew_sample[_ew_sample["Group"]=="Comparison"])), random_state=42)
_ew_plot = pd.concat([_ew_top, _ew_bad], ignore_index=True)

fig_ew_map = px.scatter_mapbox(
    _ew_plot, lat="plat", lon="plon", color="Group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison": "#ef4444"},
    zoom=9.8, center={"lat": 51.505, "lon": -0.09},
    mapbox_style="carto-darkmatter",
    opacity=0.6, height=420,
)
fig_ew_map.update_traces(marker_size=5)
fig_ew_map.update_layout(
    paper_bgcolor="#0f172a", margin=dict(l=0, r=0, t=0, b=0),
    legend=dict(orientation="h", y=0.02, x=0.01, font=dict(color="#e2e8f0", size=12),
                bgcolor="rgba(15,23,42,0.7)"),
    mapbox_layers=[dict(
        type="line",
        coordinates=[[[_WEST_LON, 51.1], [_WEST_LON, 51.9]]],
        color="#f59e0b", opacity=1.0, line=dict(width=3),
    )],
)
_add_image(s, _fig_to_img(fig_ew_map, w=1300, h=420),
           Inches(0.4), Inches(2.55), Inches(12.5), Inches(4.6))
_add_text(s, "Yellow line = Charing Cross boundary (−0.12°)  ·  Green = top drivers  ·  Red = comparison",
          Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.25),
          size=10, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5: INNER / OUTER ZONE BREAKDOWN
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Breaking London Down: Where Wait Times Stack Up",
              "Four sub-regions tell four different stories. "
              "Inner East is survivable. Outer East is where shifts go to die.")

_sub_order = ["Outer West", "Inner West", "Inner East", "Outer East"]
_sub_colors = {"Outer West": "#3b82f6", "Inner West": "#22c55e",
               "Inner East": "#f59e0b", "Outer East": "#ef4444"}

def _subregion_stats(df, label):
    rows = []
    for sub in _sub_order:
        s_df = df[df["subregion"] == sub]
        rows.append({
            "Sub-region": sub, "Group": label,
            "Avg fare £":   round(s_df["fare"].mean(), 2) if len(s_df) else 0,
            "Trip share %": round(len(s_df) / max(len(df), 1) * 100, 1),
            "Sub-£10 %":    round((s_df["fare"] < 10).mean() * 100, 1) if len(s_df) else 0,
        })
    return rows

_sub_rows = _subregion_stats(acc_top, "Top drivers") + _subregion_stats(acc_bad, "Comparison")
_sub_df = pd.DataFrame(_sub_rows)

fig_sub_fare = px.bar(
    _sub_df, x="Sub-region", y="Avg fare £", color="Group", barmode="group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison": "#ef4444"},
    category_orders={"Sub-region": _sub_order},
    text_auto=".2f", title="Avg accepted fare by sub-region",
    height=310,
)
fig_sub_fare.update_traces(texttemplate="£%{text}", textposition="outside")
fig_sub_fare.update_layout(legend=dict(orientation="h", y=1.08),
                           margin=dict(l=50, r=30, t=55, b=40))
_add_image(s, _fig_to_img(fig_sub_fare, w=760, h=310), Inches(0.4), Inches(1.3), Inches(6.5), Inches(3.6))

fig_sub_noise = px.bar(
    _sub_df, x="Sub-region", y="Sub-£10 %", color="Group", barmode="group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison": "#ef4444"},
    category_orders={"Sub-region": _sub_order},
    text_auto=".1f", title="Sub-£10 trip share (noise) by sub-region",
    height=310,
)
fig_sub_noise.update_traces(texttemplate="%{text}%", textposition="outside")
fig_sub_noise.update_layout(legend=dict(orientation="h", y=1.08),
                             margin=dict(l=50, r=30, t=55, b=40))
_add_image(s, _fig_to_img(fig_sub_noise, w=760, h=310), Inches(7.0), Inches(1.3), Inches(6.0), Inches(3.6))

# Sub-region boundary callout
_sub_ref = {sub: _sub_df[(_sub_df["Sub-region"]==sub) & (_sub_df["Group"]=="Top drivers")]["Avg fare £"].values for sub in _sub_order}
_callout(s,
    "SUB-REGION BOUNDARIES\n"
    "Outer West (lon < −0.25): Heathrow corridor, Richmond, Chiswick — long-haul heavy, "
    f"avg fare £{_sub_ref['Outer West'][0]:.2f} for top drivers.\n"
    "Inner West (−0.25 to −0.12): Kensington, Chelsea, Mayfair — premium short/medium trips.\n"
    "Inner East (−0.12 to 0.0): City of London, Shoreditch — high noise, high ceiling. Yousuf's territory.\n"
    "Outer East (lon > 0.0): East Ham, Walthamstow, Enfield — where pings dry up and fares crater. "
    "Bad drivers spend disproportionate time here and accept the scraps.",
    Inches(0.4), Inches(5.1), Inches(12.5), Inches(1.8), border_color=C_AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6: SAME STREETS, COMPLETELY DIFFERENT OUTCOMES
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Same Streets. Completely Different Outcomes.",
              "The east isn't the problem. What drivers accept in the east is the problem.")

# Yousuf vs Bartley fare scatter
_bartley_trips = acc_bad[acc_bad["dim_driver_id"]==82].copy()
_compare = pd.concat([
    acc_yousuf.assign(driver="Mohamed Yousuf (Cat A — selective)"),
    _bartley_trips.assign(driver="Aaron Bartley (Cat D — noise acceptor)"),
], ignore_index=True)
_compare["fare"] = pd.to_numeric(_compare["trip_price_in_pound"], errors="coerce")
_compare = _compare.dropna(subset=["fare","plon"])

fig_scatter = px.scatter(
    _compare, x="plon", y="fare", color="driver",
    color_discrete_map={
        "Mohamed Yousuf (Cat A — selective)": "#22c55e",
        "Aaron Bartley (Cat D — noise acceptor)": "#ef4444",
    },
    labels={"plon": "Pickup longitude", "fare": "Accepted fare (£)"},
    title="City of London corridors — accepted fare by longitude",
    height=380, opacity=0.55,
)
fig_scatter.add_vline(x=_WEST_LON, line_dash="dash", line_color="#f59e0b",
                      annotation_text="West boundary")
fig_scatter.add_hline(y=10, line_dash="dot", line_color="#94a3b8",
                      annotation_text="£10 floor", annotation_position="top right")
fig_scatter.update_layout(legend=dict(orientation="h", y=1.05))
_add_image(s, _fig_to_img(fig_scatter, w=820, h=380), Inches(0.5), Inches(1.4), Inches(7.8), Inches(4.5))

# Callout cards
_y_sub10 = round((_compare[_compare["driver"].str.contains("Yousuf")]["fare"] < 10).mean() * 100, 0)
_b_sub10 = round((_compare[_compare["driver"].str.contains("Bartley")]["fare"] < 10).mean() * 100, 0)
_y_avg   = round(_compare[_compare["driver"].str.contains("Yousuf")]["fare"].mean(), 2)
_b_avg   = round(_compare[_compare["driver"].str.contains("Bartley")]["fare"].mean(), 2)

_add_rect(s, Inches(8.7), Inches(1.4), Inches(4.3), Inches(2.1), fill=RGBColor(0x05,0x1e,0x0a))
_add_text(s, "YOUSUF — City of London", Inches(8.85), Inches(1.45), Inches(4.0), Inches(0.4),
          size=11, bold=True, color=C_GREEN)
_add_text(s, f"Avg fare:      £{_y_avg:.2f}\nSub-£10 trips: {_y_sub10:.0f}%\nStrategy:      10% acceptance — waits for quality",
          Inches(8.85), Inches(1.85), Inches(4.0), Inches(1.3), size=12, color=C_WHITE)

_add_rect(s, Inches(8.7), Inches(3.7), Inches(4.3), Inches(2.1), fill=RGBColor(0x1f,0x05,0x05))
_add_text(s, "BARTLEY — Same Streets", Inches(8.85), Inches(3.75), Inches(4.0), Inches(0.4),
          size=11, bold=True, color=C_RED)
_add_text(s, f"Avg fare:      £{_b_avg:.2f}\nSub-£10 trips: {_b_sub10:.0f}%\nStrategy:      ~88% acceptance — takes all pings",
          Inches(8.85), Inches(4.15), Inches(4.0), Inches(1.3), size=12, color=C_WHITE)

_callout(s,
    "The east isn't the problem. The noise is. Good drivers operate in inner east with a strict fare+distance filter. "
    "Bad drivers accept the same streets' cheap sub-£10 hops — then wonder why their shift ends at £18/hr.",
    Inches(0.5), Inches(6.1), Inches(12.5), Inches(0.9), border_color=C_AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6: ZONE EFFICIENCY — NUMBERS DON'T LIE
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Zone Efficiency: Z1 Earns Nearly Double Z3 Per Real Hour",
              "True RPH = fare ÷ (inter-trip gap + pickup time + ride time). This is the number that matters.")

zone_eff = pd.DataFrame([
    {"Zone": "Zone 1", "True RPH £": 21.0, "Avg wait (min)":  7},
    {"Zone": "Zone 2", "True RPH £": 19.5, "Avg wait (min)": 12},
    {"Zone": "Zone 3", "True RPH £": 17.5, "Avg wait (min)": 32},
    {"Zone": "Zone 4", "True RPH £": 16.0, "Avg wait (min)": 28},
    {"Zone": "Zone 5", "True RPH £": 15.0, "Avg wait (min)": 35},
    {"Zone": "Zone 6", "True RPH £": 23.0, "Avg wait (min)": 18},
])

fig_z_rph = px.bar(
    zone_eff, x="Zone", y="True RPH £",
    color="True RPH £", color_continuous_scale="RdYlGn",
    text="True RPH £", title="True RPH by Pickup Zone (incl. wait time)", height=370,
)
fig_z_rph.update_traces(texttemplate="£%{text:.0f}", textposition="outside")
fig_z_rph.update_layout(coloraxis_showscale=False, yaxis=dict(range=[0,28]))
_add_image(s, _fig_to_img(fig_z_rph, w=720, h=370), Inches(0.5), Inches(1.3), Inches(6.2), Inches(4.2))

fig_z_wait = px.bar(
    zone_eff, x="Zone", y="Avg wait (min)",
    color="Avg wait (min)", color_continuous_scale="RdYlGn_r",
    text="Avg wait (min)", title="Avg Inter-trip Wait Time by Zone", height=370,
)
fig_z_wait.update_traces(texttemplate="%{text} min", textposition="outside")
fig_z_wait.update_layout(coloraxis_showscale=False, yaxis=dict(range=[0,42]))
_add_image(s, _fig_to_img(fig_z_wait, w=720, h=370), Inches(6.9), Inches(1.3), Inches(6.1), Inches(4.2))

_callout(s,
    "Zone 3 daytime is the trap: £17.50/hr true RPH with a 32-minute average wait for the next ping. "
    "Zone 6 looks best but drops passengers far from the city — return positioning eats the gain. "
    "Zone 1 is the consistent winner: short waits, high-density pings, £21/hr true cycle.\n\n"
    "Note: Zone 3 at night/early morning (00:00–09:00) reaches £32–43/hr — "
    "the problem is specifically daytime Z3.",
    Inches(0.5), Inches(5.5), Inches(12.5), Inches(1.6), border_color=C_BLUE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7: ZONE FLOW — THE Z3 TRAP
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "The Zone 3 Gravity Well: Once You're In, You Stay In",
              "Bad drivers are not just visiting Zone 3 more — they're chaining trips inside it.")

_metric_card(s, "Z1 pickup share — top drivers",  f"{z1_top:.1f}%", Inches(0.5), Inches(1.4), val_color=C_GREEN)
_metric_card(s, "Z1 pickup share — bad drivers",  f"{z1_bad:.1f}%", Inches(3.1), Inches(1.4), val_color=C_RED,
             delta=f"{z1_bad - z1_top:+.1f}% gap", delta_color=C_RED)
_metric_card(s, "Z3 pickup share — top drivers",  f"{z3_top:.1f}%", Inches(5.7), Inches(1.4), val_color=C_GREEN)
_metric_card(s, "Z3 pickup share — bad drivers",  f"{z3_bad:.1f}%", Inches(8.3), Inches(1.4), val_color=C_RED,
             delta=f"{z3_bad - z3_top:+.1f}% more in Z3", delta_color=C_RED)
_metric_card(s, "Z3→Z3 chain — bad drivers",      f"{chain_bad:.1f}%", Inches(10.9), Inches(1.4), val_color=C_RED)

zone_cmp = pd.DataFrame([
    {"Metric": "Z1 pickups",   "Top drivers": z1_top,    "Comparison": z1_bad},
    {"Metric": "Z3 pickups",   "Top drivers": z3_top,    "Comparison": z3_bad},
    {"Metric": "Z3→Z3 chain",  "Top drivers": chain_top, "Comparison": chain_bad},
])
fig_z3cmp = px.bar(
    zone_cmp.melt(id_vars="Metric", var_name="Group", value_name="% of trips"),
    x="Metric", y="% of trips", color="Group", barmode="group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison": "#ef4444"},
    title="Zone 1 vs Zone 3 activity — Top vs Comparison (%)",
    text_auto=".1f", height=350,
)
fig_z3cmp.update_layout(legend=dict(orientation="h", y=1.05))
_add_image(s, _fig_to_img(fig_z3cmp, w=760, h=350), Inches(0.5), Inches(2.75), Inches(7.4), Inches(4.0))

_callout(s,
    "Zone 3 is a gravity well. Drivers picking up in Z3 receive the next ping from a nearby Z3 passenger "
    "— because the ping system rewards proximity. To escape, a driver must deliberately decline Z3 pings "
    "and wait for one that pulls them toward Z1/Z2. Bad drivers never make that move.\n\n"
    f"Bad drivers run {chain_bad:.0f}% of trips as Z3→Z3 chains. "
    f"Top drivers: {chain_top:.0f}%. Every Z3→Z3 chain is another 32 minutes of below-average RPH locked in.",
    Inches(8.1), Inches(2.75), Inches(5.0), Inches(4.0), border_color=C_RED)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8: ZONE 3 NUANCE — DAYTIME TRAP, NIGHT OPPORTUNITY
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Zone 3 Is Not Universally Bad — Timing Is Everything",
              "The trap is daytime Z3. Night and early morning Z3 is a different story entirely.")

z3_time = pd.DataFrame([
    {"Period": "Night (00–06)",    "Raw RPH £": 36, "Verdict": "Worth it"},
    {"Period": "Morning (06–09)",  "Raw RPH £": 34, "Verdict": "Worth it"},
    {"Period": "Day (09–17)",      "Raw RPH £": 27, "Verdict": "Avoid"},
    {"Period": "Evening (17–23)",  "Raw RPH £": 29, "Verdict": "Borderline"},
])
fig_z3t = px.bar(
    z3_time, x="Period", y="Raw RPH £", color="Verdict",
    color_discrete_map={"Worth it": "#22c55e", "Avoid": "#ef4444", "Borderline": "#f59e0b"},
    text="Raw RPH £", title="Zone 3 Raw RPH by time of day",
    height=360,
)
fig_z3t.add_hline(y=21, line_dash="dash", line_color="#3b82f6",
                  annotation_text="Zone 1 avg £21/hr (with wait)", annotation_position="top right")
fig_z3t.update_traces(texttemplate="£%{text}", textposition="outside")
fig_z3t.update_layout(legend=dict(orientation="h", y=1.05))
_add_image(s, _fig_to_img(fig_z3t, w=780, h=360), Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.5))

_callout(s,
    "NIGHT SHIFT IN Z3 (00:00–09:00)\n"
    "Longer trips (avg 11 miles) heading into the city. Fewer competing drivers. "
    "RPH reaches £36/hr — above the Z1 average. Best Z3 areas at night: Ealing, Brent, Wandsworth.\n\n"
    "DAYTIME Z3 (09:00–17:00)\n"
    "Short local hops, 32-minute avg wait, sub-£18/hr true RPH. "
    "Worst areas: Enfield, East Ham, Walthamstow.\n\n"
    "The rule: if you are in Z3 after 09:00, decline until you get a Z1/Z2 pull or an airport.",
    Inches(8.3), Inches(1.4), Inches(4.8), Inches(5.5), border_color=C_AMBER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9: GAPS — STRATEGIC WAIT vs DEAD TIME
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Not All Gaps Are Bad — Location Determines Everything",
              "A 30-minute gap in Zone 1 is a strategic wait. The same gap in outer east is dead time with zero pings incoming.")

def _gap_buckets(df):
    cuts  = [0, 10, 25, 60, 180]
    labs  = ["<10 min", "10–25 min", "25–60 min", ">60 min"]
    df    = df.copy()
    df["bucket"] = pd.cut(df["gap_min"], bins=cuts, labels=labs, right=True)
    return df["bucket"].value_counts(normalize=True).reindex(labs).fillna(0) * 100

gb_top = _gap_buckets(gaps_top_df)
gb_bad = _gap_buckets(gaps_bad_df)
gap_df = pd.DataFrame({
    "Gap bucket": gb_top.index,
    "Top drivers": gb_top.values,
    "Comparison":  gb_bad.values,
})
fig_gap = px.bar(
    gap_df.melt(id_vars="Gap bucket", var_name="Group", value_name="% of gaps"),
    x="Gap bucket", y="% of gaps", color="Group", barmode="group",
    color_discrete_map={"Top drivers": "#22c55e", "Comparison": "#ef4444"},
    title="Gap duration distribution — Top vs Comparison (%)",
    text_auto=".1f", height=350,
)
fig_gap.update_layout(legend=dict(orientation="h", y=1.05))
_add_image(s, _fig_to_img(fig_gap, w=780, h=350), Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.5))

_metric_card(s, "Top drivers — avg gap",  f"{gap_top_avg:.0f} min", Inches(8.4), Inches(1.4),
             val_color=C_GREEN, w=Inches(2.3))
_metric_card(s, "Bad drivers — avg gap",  f"{gap_bad_avg:.0f} min", Inches(11.0), Inches(1.4),
             val_color=C_RED, w=Inches(2.0))

_callout(s,
    "WHY A LONGER GAP CAN BE A GOOD SIGN\n"
    "Top drivers sometimes wait 25–60 minutes — but that wait is in a high-ping corridor (Z1/Z2). "
    "The trip that follows averages significantly more than a trip taken immediately from a bad position.\n\n"
    "WHY THE SAME GAP FOR BAD DRIVERS IS JUST DEAD TIME\n"
    "Drivers stranded in outer east/Zone 3 after a dropoff receive fewer pings per minute. "
    "The 32-minute average Z3 wait produces no revenue — it's just the system struggling to find them a next trip.",
    Inches(8.4), Inches(2.75), Inches(4.6), Inches(3.3), border_color=C_BLUE)

_add_text(s, f"Avg gap top: {gap_top_avg:.0f} min  ·  Avg gap bad: {gap_bad_avg:.0f} min  ·  "
             "Gap length alone does not explain RPH — zone at time of gap does.",
          Inches(0.5), Inches(6.8), Inches(12.5), Inches(0.35),
          size=10, color=C_MUTED, italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10: WHAT GOOD DRIVERS REJECT — Z1 SELECTIVITY
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Good Drivers Reject the Most — Especially in Zone 1",
              "Counterintuitive finding: top drivers decline a higher share of Zone 1 pings than any other zone. "
              "They are not avoiding Z1 — they are waiting for the right Z1 trip.")

accept_zone = pd.DataFrame({
    "Zone":     ["Zone 1","Zone 2","Zone 3","Zone 4","Zone 5","Zone 6"],
    "Accepted": [27.8, 26.6, 24.3, 10.1, 5.6, 5.6],
    "Declined": [34.9, 28.2, 23.1,  7.7, 3.3, 2.8],
})
fig_acc_zone = px.bar(
    accept_zone.melt(id_vars="Zone", var_name="Outcome", value_name="% of pings"),
    x="Zone", y="% of pings", color="Outcome", barmode="group",
    color_discrete_map={"Accepted": "#22c55e", "Declined": "#ef4444"},
    title="Accepted vs Declined pings by pickup zone — top 10 drivers",
    text_auto=".1f", height=350,
)
fig_acc_zone.update_layout(legend=dict(orientation="h", y=1.05))
_add_image(s, _fig_to_img(fig_acc_zone, w=780, h=350), Inches(0.5), Inches(1.4), Inches(7.5), Inches(4.5))

# Fare distribution
_fares_top = pd.to_numeric(acc_top["trip_price_in_pound"], errors="coerce").dropna()
_fares_bad = pd.to_numeric(acc_bad["trip_price_in_pound"], errors="coerce").dropna()
combined_fare = pd.DataFrame({
    "Fare £": pd.concat([_fares_top, _fares_bad], ignore_index=True),
    "Group":  ["Top drivers"] * len(_fares_top) + ["Comparison"] * len(_fares_bad),
})
combined_fare = combined_fare[combined_fare["Fare £"].between(5, 80)]
top_med = combined_fare[combined_fare["Group"]=="Top drivers"]["Fare £"].median()
bad_med = combined_fare[combined_fare["Group"]=="Comparison"]["Fare £"].median()

fig_fd = px.histogram(
    combined_fare, x="Fare £", color="Group",
    barmode="overlay", nbins=35, opacity=0.7,
    color_discrete_map={"Top drivers":"#22c55e","Comparison":"#ef4444"},
    title="Accepted fare distribution", height=310,
)
fig_fd.add_vline(x=top_med, line_dash="dash", line_color="#22c55e",
                 annotation_text=f"Top £{top_med:.2f}")
fig_fd.add_vline(x=bad_med, line_dash="dash", line_color="#ef4444",
                 annotation_text=f"Bad £{bad_med:.2f}", annotation_position="top left")
fig_fd.update_layout(legend=dict(orientation="h", y=1.05))
_add_image(s, _fig_to_img(fig_fd, w=580, h=380), Inches(8.3), Inches(1.4), Inches(4.8), Inches(4.5))

_callout(s,
    "The Z1 selectivity insight: top drivers park in Z1/Z2 and decline the short-hop, low-fare pings "
    "until a cross-zone trip appears. They are using Z1's high ping density as a waiting room — "
    "not as a mandate to accept everything. This is the exact pattern the model needs to encode.",
    Inches(0.5), Inches(6.1), Inches(12.5), Inches(0.9), border_color=C_GREEN)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11: ACCEPTANCE RATE — THE FLEET IS OVERSHOOTING IN BOTH DIRECTIONS
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "Acceptance Rate: Lower Acceptance, Higher Earnings — The Data Proves It",
              "April: 53% acceptance, £20.74/hr.  May: 43% acceptance, £22.44/hr. "
              "The fleet over-corrected but the direction was right. Target is 50% — where top drivers naturally land.")

# Left: acceptance rate history bar
acc_hist = pd.DataFrame({
    "Period":          ["Historical\nfleet", f"Top 10\ndrivers\n({top_acc:.0f}%)", "April\n2026", "May\n2026", "Bolt\ntarget"],
    "Acceptance rate": [57.2,               top_acc,                               53.0,          43.0,        50.0],
    "Color":           ["#64748b",           "#22c55e",                             "#3b82f6",     "#ef4444",   "#f59e0b"],
})
fig_acc_hist = go.Figure()
for _, row in acc_hist.iterrows():
    fig_acc_hist.add_bar(x=[row["Period"]], y=[row["Acceptance rate"]],
                         name=row["Period"], marker_color=row["Color"],
                         text=f"{row['Acceptance rate']:.0f}%", textposition="outside")
fig_acc_hist.add_hline(y=50, line_dash="dash", line_color="#f59e0b",
                       annotation_text="Bolt target 50%", annotation_position="top right")
fig_acc_hist.update_layout(showlegend=False, yaxis_title="Acceptance rate (%)",
                           height=330, title="Acceptance rate — fleet history vs target",
                           yaxis=dict(range=[0,72]), margin=dict(l=50,r=30,t=50,b=40))
_add_image(s, _fig_to_img(fig_acc_hist, w=680, h=330), Inches(0.4), Inches(1.3), Inches(6.0), Inches(3.8))

# Right: April vs May dual-axis — acceptance rate vs RPH (the KEY new repo finding)
fig_apr_may = go.Figure()
months   = ["April 2026", "May 2026"]
acc_vals = [53.0, 43.0]
rph_vals = [20.74, 22.44]

fig_apr_may.add_bar(x=months, y=acc_vals, name="Acceptance rate (%)",
                    marker_color=["#3b82f6","#ef4444"],
                    text=[f"{v:.0f}%" for v in acc_vals], textposition="outside",
                    yaxis="y1")
fig_apr_may.add_scatter(x=months, y=rph_vals, name="Avg RPH (£)",
                        mode="lines+markers+text",
                        line=dict(color="#f59e0b", width=3),
                        marker=dict(size=12, color="#f59e0b"),
                        text=[f"£{v:.2f}" for v in rph_vals],
                        textposition="top center",
                        textfont=dict(color="#f59e0b", size=14),
                        yaxis="y2")
fig_apr_may.update_layout(
    title="April vs May: Acceptance ↓, RPH ↑",
    yaxis=dict(title="Acceptance rate (%)", range=[0,70], color="#94a3b8"),
    yaxis2=dict(title="RPH (£)", range=[18,26], overlaying="y", side="right",
                color="#f59e0b", showgrid=False),
    legend=dict(orientation="h", y=1.1),
    height=330, margin=dict(l=50,r=60,t=55,b=40),
    barmode="group",
)
_add_image(s, _fig_to_img(fig_apr_may, w=680, h=330), Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.8))

# Metric strip
_metric_card(s, "April acceptance",  "53%",     Inches(0.4),  Inches(5.35), val_color=C_BLUE,  w=Inches(2.5))
_metric_card(s, "April RPH",         "£20.74",  Inches(3.1),  Inches(5.35), val_color=C_BLUE,  w=Inches(2.5))
_metric_card(s, "May acceptance",    "43%",     Inches(5.8),  Inches(5.35), val_color=C_RED,   w=Inches(2.5))
_metric_card(s, "May RPH",           "£22.44",  Inches(8.5),  Inches(5.35), val_color=C_GREEN, w=Inches(2.5),
             delta="+£1.70/hr with 10pp fewer accepts", delta_color=C_GREEN)

_callout(s,
    "AIRPORT DRAG — hidden suppressor of the current 43%: Drivers queuing at Heathrow cannot accept pings. "
    "Every missed ping during queue time counts as a decline. This alone accounts for an estimated 1–2pp "
    "of the 43% figure and cannot be fixed through driver behaviour — Bolt must exclude queue windows from the KPI.",
    Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.7), border_color=C_RED)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12: THE MODEL — WHAT PARAMETERS IT NEEDS
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "The Model: Good Drivers Have Already Built the Decision Tree",
              "Every accept/reject decision from our top 10 is a labelled training example. "
              "The model already exists — this analysis tells us what parameters to sharpen.")

params = [
    ("PICKUP ZONE",          "Z1/Z2 vs Z3+ at time of ping",
     "Z1/Z2 = accept unless sub-2-mile. Z3 daytime = decline unless airport/long-haul.",
     C_GREEN),
    ("TIME OF DAY",          "Peak (07–09, 17–20) vs off-peak vs night",
     "Peak: relax fare floor. Off-peak: tighten. Night Z3: reopen.",
     C_BLUE),
    ("EAST / WEST",          "Pickup longitude vs −0.12° boundary",
     "West: accept wider fare range. East inner: high threshold. East outer: only long-haul.",
     C_ACCENT),
    ("ESTIMATED DISTANCE",   "Trip length, not just fare",
     "Sub-2-mile = decline regardless of fare. 2–5 mile = evaluate zone+fare.",
     C_AMBER),
    ("DROPOFF ZONE",         "Will this trip strand the driver?",
     "Most important addition: score not just the trip but where it leaves the driver.",
     RGBColor(0xa8,0x5e,0xf7)),
    ("CURRENT GAP STATE",    "How long has the driver been waiting?",
     "Driver in 30-min gap in Z1 = wait longer. Driver in Z3 dead time = accept almost anything to escape.",
     C_MUTED),
]

for i, (title, subtitle, detail, col) in enumerate(params):
    col_i = i % 3
    row_i = i // 3
    lft = Inches(0.4) + col_i * Inches(4.3)
    top = Inches(1.5) + row_i * Inches(2.6)
    _add_rect(s, lft, top, Inches(4.0), Inches(2.3), fill=C_CARD)
    _add_rect(s, lft, top, Inches(4.0), Inches(0.08), fill=col)
    _add_text(s, title, lft + Inches(0.12), top + Inches(0.15), Inches(3.8), Inches(0.35),
              size=11, bold=True, color=col)
    _add_text(s, subtitle, lft + Inches(0.12), top + Inches(0.5), Inches(3.8), Inches(0.35),
              size=10, color=C_MUTED, italic=True)
    _add_text(s, detail, lft + Inches(0.12), top + Inches(0.9), Inches(3.8), Inches(1.2),
              size=11, color=C_WHITE)

_add_text(s, "★  Dropoff zone scoring is the highest-priority addition — shift the model from single-trip evaluation to sequence optimisation.",
          Inches(0.4), Inches(6.9), Inches(12.5), Inches(0.4),
          size=11, color=RGBColor(0xa8,0x5e,0xf7), italic=True)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13: THE ASK — WHAT MANAGEMENT NEEDS TO APPROVE
# ════════════════════════════════════════════════════════════════════════════
s = _blank_slide(prs)
_slide_header(s, "What We Need to Approve",
              "Five operational changes. One model update. One message to Bolt.")

actions = [
    ("1", "POSITIONING GUIDELINES",
     "Publish Z1/Z2 dwell guidance to all drivers. Daytime Z3 avoidance protocol. "
     "Simple: 'If you're in Z3 after 9am, decline until you get a Z1 pull.'",
     C_GREEN),
    ("2", "PEAK-HOUR FARE FLOOR",
     "Lower minimum fare threshold during 07–09 and 17–20 for all drivers. "
     "Data shows sub-£12 trips in those windows are RPH-neutral due to short gaps.",
     C_BLUE),
    ("3", "DISTANCE-BASED FILTER",
     "Replace fare-floor guidance with distance-floor guidance. "
     "Sub-2-mile trips are the RPH killers — not sub-£12 trips.",
     C_ACCENT),
    ("4", "WEST SHIFT INCENTIVE",
     "Review whether dispatcher logic can weight westward positioning. "
     "Bad drivers operating outer east get genuinely bad pings — structural fix, not driver behaviour fix.",
     C_AMBER),
    ("5", "MODEL PARAMETER UPDATE",
     "Add dropoff zone scoring + time-of-day weights to the existing model. "
     "This makes recommendations sequence-aware, not single-trip-aware.",
     RGBColor(0xa8,0x5e,0xf7)),
    ("B", "MESSAGE TO BOLT",
     "Heathrow queue dead time suppresses acceptance rate — drivers miss pings while queuing. "
     "Request exclusion of queue windows from acceptance rate KPI calculation. "
     "This alone may account for 1–2pp of the current 43% figure.",
     C_RED),
]

for i, (num, title, detail, col) in enumerate(actions):
    col_i = i % 2
    row_i = i // 2
    lft = Inches(0.4)  + col_i * Inches(6.5)
    top = Inches(1.45) + row_i * Inches(1.85)
    _add_rect(s, lft, top, Inches(6.1), Inches(1.65), fill=C_CARD)
    _add_rect(s, lft, top, Inches(0.45), Inches(1.65), fill=col)
    _add_text(s, num, lft + Inches(0.08), top + Inches(0.55), Inches(0.3), Inches(0.5),
              size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, title, lft + Inches(0.55), top + Inches(0.08), Inches(5.4), Inches(0.4),
              size=12, bold=True, color=col)
    _add_text(s, detail, lft + Inches(0.55), top + Inches(0.5), Inches(5.4), Inches(1.05),
              size=11, color=C_WHITE)

# ── Save ─────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "odysse_management_deck.pptx")
prs.save(out_path)
print(f"\nSaved: {out_path}")
print(f"Slides: {len(prs.slides)}")
